from pptx import Presentation
from pptx.util import Inches, Pt

# Ustvari novo prezentacijo
prs = Presentation()

# Funkcija za dodajanje slida z naslovom in vsebino (Layout 1: Title and Content)
def add_content_slide(prs, title_text, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Nastavi naslov
    title = slide.shapes.title
    title.text = title_text
    
    # Nastavi vsebino kot alinee
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Počisti prednastavljeno besedilo
    
    for point in content_list:
        p = tf.add_paragraph()
        p.text = point
        # Nastavi velikost pisave (za boljši videz)
        p.font.size = Pt(20) 
        p.level = 0

# --- SLIDE 1: NASLOVNICA (GOVOREC 1) ---
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ATHLETIQ"
subtitle.text = "AI VEČ-ŠPORTNI TRENER\n\nTrening • Prehrana • Regeneracija\nPersonalizirano z umetno inteligenco"

# --- SLIDE 2: PROBLEM (GOVOREC 1) ---
add_content_slide(prs, "IZZIV: Fragmentirano in drago trenerstvo", [
    "Generični programi: Treningi niso prilagojeni posamezniku (igre z žogo, maraton...).",
    "Stroški: Osebni trenerji so za večino športnikov predragi ali nedostopni.",
    "Fragmentacija: Trening, prehrana in regeneracija so ločeni sistemi/aplikacije.",
    "Neizkoriščenost: Pametne naprave zbirajo podatke, ki ne vodijo do dejanskega napredka."
])

# --- SLIDE 3: REŠITEV (GOVOREC 1) ---
add_content_slide(prs, "REŠITEV: Tvoj osebni AI trener 24/7", [
    "🧠 En digitalni trener, ki se popolnoma prilagaja posamezniku.",
    "🔄 Umetna inteligenca uporablja podatke za dinamično prilagajanje v realnem času.",
    "🔗 Povezava: Povezuje in optimizira trening, prehrano in regeneracijo.",
    "🚀 Napredek: Podatke uporabnika pretvori v rezultate in napredek."
])

# --- SLIDE 4: TRG IN CILJNA PUBLIKA (GOVOREC 2) ---
add_content_slide(prs, "TRG IN CILJNA PUBLIKA", [
    "📈 Tržna priložnost: Velik in rastoč globalni trg fitnes aplikacij in wearables.",
    "🎯 Ciljna publika: Aktivni rekreativni in polprofesionalni športniki (18-45 let).",
    "🧑‍💻 Profili: Uporabniki, ki aktivno uporabljajo tehnologijo in želijo strukturo in napredek.",
    "✅ Priložnost: Generične aplikacije ne pokrivajo nišnih in ekipnih športov."
])

# --- SLIDE 5: PRODUKTNA VIZIJA IN MVP (GOVOREC 4) ---
add_content_slide(prs, "PRODUKTNA VIZIJA IN MVP", [
    "Vizija: Ustvariti digitalnega trenerja, ki se sčasoma uči in prilagaja telesu uporabnika.",
    "UX/UI: Uporabniška izkušnja mora biti preprosta, hitra in motivacijska.",
    "MVP vključuje: Personalizirane treninge za vsak šport (generativni AI).",
    "MVP vključuje: Osnovne prehranske in regeneracijske predloge. Povezava z Apple Health/Garmin."
])

# --- SLIDE 6: KONKURENČNA PREDNOST IN POZICIONIRANJE (GOVOREC 2 + 3) ---
add_content_slide(prs, "KONKURENČNA PREDNOST", [
    "🚀 Prednost: Neskončna razširljivost športov – ni ročnega načrtovanja treningov.",
    "💡 Vloga: Nismo 'tracker', ampak smo 'coach'.",
    "🎯 **Positioning Statement:** AthletiQ je AI trener, ki povezuje trening, prehrano in regeneracijo v en pameten sistem."
])

# --- SLIDE 7: BLAGOVNA ZNAMKA IN GO-TO-MARKET (GOVOREC 3) ---
add_content_slide(prs, "POZICIONIRANJE IN ZNAMKA", [
    "One-liner: **Train smarter. Recover better.**",
    "Blagovna znamka: Strokoven, samozavesten in usmerjen v dolgoročni napredek.",
    "Strategija: Fokus na kanale, kjer je naša ciljna publika že prisotna."
])

# --- SLIDE 8: KANALI IN RAST (GOVOREC 3) ---
add_content_slide(prs, "KANALI IN NARAVNI 'GROWTH LOOP'", [
    "Organski kanali: Instagram in TikTok (vsebina, izzivi, rezultati uporabnikov).",
    "Plačani kanali: Meta in Google oglasi, usmerjeni na nišne športne skupine.",
    "Growth Loop: Če uporabnik vidi napredek, aplikacijo deli – kar ustvarja naravni cikel rasti."
])

# --- SLIDE 9: POSLOVNI MODEL IN FINANCE ---
add_content_slide(prs, "POSLOVNI MODEL IN TRAKCIJA", [
    "Model: Freemium (brezplačni osnovni treningi + oglasi) in Premium.",
    "Premium (20 €/mesec): Popolna personalizacija, AI analiza obrokov.",
    "Projekcija Leto 1: 100.000 uporabnikov / 720.000 € ARR (Zagon).",
    "Projekcija Leto 3: 1.000.000 uporabnikov / 12.000.000 € ARR (Globalna širitev)."
])

# --- SLIDE 10: INVESTICIJA (ASK) ---
add_content_slide(prs, "PRILOŽNOST ZA INVESTICIJO", [
    "💶 Iščemo: 750.000 € (Seed runda).",
    "🏗 Uporaba sredstev: 40% AI & infrastruktura, 30% Ekipa, 20% Marketing, 10% Operativa.",
    "Cilj: Lansiranje v EU in ZDA v 12 mesecih."
])

# --- SLIDE 11: ZAKLJUČEK (GOVOREC 4) ---
add_content_slide(prs, "ZAKLJUČEK IN VIZIJA", [
    "🌟 Verjamemo: Vsak športnik si zasluži elitnega trenerja.",
    "🚀 Misija: Dvigniti nivo rekreativnega športa z uporabo AI.",
    "❤️ Vabimo vas, da se pridružite oblikovanju prihodnosti športne uspešnosti."
])

# Shrani prezentacijo
prs.save('AthletiQ_Revised_Pitch_Deck2.pptx')
print("Posodobljena prezentacija ustvarjena: AthletiQ_Revised_Pitch_Deck.pptx")